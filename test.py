from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
# 读入已存在的PPT文件
prs = Presentation()
slides = prs.slides
slide = prs.slides.add_slide(prs.slide_layouts[6])
# 添加表格
x,y,cx,cy = Inches(2),Inches(2),Inches(15),Inches(2)
shape = slide.shapes.add_table(2,16,x,y,cx,cy)#3*3的表格  参数rows,cols,left,top,width,height
# shape.has_table
table = shape.table
# 单元格中添加值
# cell = table.
# cell.text = 'Apple'
R,G,B = 237, 69, 87
N_R,N_G,N_B =  250,212,232
P_R,P_G,P_B = 217, 244, 245
for row in range(len(table.rows)):
    for col in range(len(table.columns)):
        if row == 0:
            table.cell(row, col).text = "第一行"
            table.cell(row, col).text_frame.paragraphs[0].font.size=Pt(13)
            table.cell(row, col).text_frame.paragraphs[0].font.name="微软雅黑"
            table.cell(row, col).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 居中
            table.cell(row, col).text_frame.paragraphs[0].font.blod = True #加粗
            table.cell(row, col).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,206,209)
            table.cell(row, col).fill.solid()
            table.cell(row, col).fill.fore_color.rgb=RGBColor(240,255,255)  # 背景颜色
            #Write rest of table entries
        else:
            table.cell(row, col).text = "其他行"
            table.cell(row, col).text_frame.paragraphs[0].font.size=Pt(11)
            table.cell(row, col).text_frame.paragraphs[0].font.name="微软雅黑"
            table.cell(row, col).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 居中
            table.cell(row, col).text_frame.paragraphs[0].font.color.rgb = RGBColor(R, G, B)
            table.cell(row, col).fill.solid()
            table.cell(row, col).fill.fore_color.rgb=RGBColor(255,255,255)
            if col==1:
                table.cell(row, col).text_frame.paragraphs[0].font.color.rgb=RGBColor(255,48,48)
            elif row%2==1 and R==237:
                table.cell(row, col).text_frame.paragraphs[0].font.color.rgb=RGBColor(28,28,28)
            else:

                table.cell(row, col).fill.transparency=200
prs.save('test01.pptx')